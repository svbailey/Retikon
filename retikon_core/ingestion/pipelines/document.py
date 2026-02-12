from __future__ import annotations

import os
import re
import uuid
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass
from datetime import datetime, timezone
from functools import lru_cache
from pathlib import Path

import fitz
import pandas as pd
from docx import Document
from pptx import Presentation

from retikon_core.config import Config
from retikon_core.embeddings import get_embedding_backend, get_text_embedder
from retikon_core.embeddings.timeout import run_inference
from retikon_core.errors import PermanentError
from retikon_core.ingestion.ocr import ocr_text_from_pdf
from retikon_core.ingestion.pipelines.embedding_utils import text_embed_batch_size
from retikon_core.ingestion.pipelines.metrics import (
    CallTracker,
    StageTimer,
    build_stage_timings,
    timed_call,
)
from retikon_core.ingestion.pipelines.types import PipelineResult
from retikon_core.ingestion.types import IngestSource
from retikon_core.storage.manifest import (
    build_manifest,
    manifest_metrics_subset,
    write_manifest,
)
from retikon_core.storage.paths import (
    edge_part_uri,
    manifest_uri,
    vertex_part_uri,
)
from retikon_core.storage.schemas import schema_for
from retikon_core.storage.writer import WriteResult, write_parquet
from retikon_core.tenancy import tenancy_fields


@dataclass(frozen=True)
class Chunk:
    index: int
    text: str
    char_start: int
    char_end: int
    token_start: int
    token_end: int
    token_count: int


def _pipeline_model() -> str:
    return os.getenv("TEXT_MODEL_NAME", "BAAI/bge-base-en-v1.5")


def _tokenizer_name() -> str:
    return os.getenv("TEXT_MODEL_NAME", "BAAI/bge-base-en-v1.5")


def _tokenizer_cache_dir() -> str | None:
    return os.getenv("MODEL_DIR")


def _doc_parquet_compression() -> str:
    value = os.getenv("DOC_PARQUET_COMPRESSION", "").strip()
    if value:
        return value
    return os.getenv("PARQUET_COMPRESSION", "zstd").strip() or "zstd"


def _doc_parquet_row_group_size() -> int | None:
    raw = os.getenv("DOC_PARQUET_ROW_GROUP_SIZE", "").strip()
    if not raw:
        raw = os.getenv("PARQUET_ROW_GROUP_SIZE", "").strip()
    if not raw:
        return None
    try:
        value = int(raw)
    except ValueError:
        return None
    return value if value > 0 else None


def _write_parquet_parallel(
    jobs: list[tuple[list[dict[str, object]], object, str]],
    *,
    compression: str,
    row_group_size: int | None,
) -> list[WriteResult]:
    if not jobs:
        return []
    if len(jobs) == 1:
        rows, schema, uri = jobs[0]
        return [
            write_parquet(
                rows,
                schema,
                uri,
                compression=compression,
                row_group_size=row_group_size,
            )
        ]
    max_workers = min(4, len(jobs))
    results: list[WriteResult | None] = [None] * len(jobs)
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_map = {}
        for idx, (rows, schema, uri) in enumerate(jobs):
            future = executor.submit(
                write_parquet,
                rows,
                schema,
                uri,
                compression,
                row_group_size,
            )
            future_map[future] = idx
        for future in as_completed(future_map):
            results[future_map[future]] = future.result()
    return [item for item in results if item is not None]


class _SimpleTokenizer:
    def __call__(
        self,
        text: str,
        *,
        return_offsets_mapping: bool = True,
        add_special_tokens: bool = False,
    ) -> dict[str, list]:
        offsets: list[tuple[int, int]] = []
        input_ids: list[int] = []
        if text:
            for idx, match in enumerate(re.finditer(r"\S+", text)):
                offsets.append((match.start(), match.end()))
                input_ids.append(idx)
        return {
            "input_ids": input_ids,
            "offset_mapping": offsets if return_offsets_mapping else [],
        }


def _use_simple_tokenizer() -> bool:
    return os.getenv("RETIKON_TOKENIZER", "").lower() in {
        "stub",
        "simple",
        "whitespace",
    }


@lru_cache(maxsize=1)
def _load_tokenizer():
    if _use_simple_tokenizer():
        return _SimpleTokenizer()

    try:
        from transformers import AutoTokenizer
    except ImportError:
        return _SimpleTokenizer()

    cache_dir = _tokenizer_cache_dir()
    return AutoTokenizer.from_pretrained(
        _tokenizer_name(),
        cache_dir=cache_dir,
        use_fast=True,
    )


def _extract_text(path: str, extension: str) -> str:
    if extension == ".pdf":
        doc = fitz.open(path)
        try:
            return "\n".join(page.get_text() for page in doc)
        finally:
            doc.close()
    if extension == ".docx":
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    if extension == ".pptx":
        pres = Presentation(path)
        parts: list[str] = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = getattr(shape, "text", "")
                    if text:
                        parts.append(text)
        return "\n".join(parts)
    if extension in (".csv", ".tsv"):
        sep = "," if extension == ".csv" else "\t"
        df = pd.read_csv(path, sep=sep)
        return _table_to_text(df)
    if extension in (".xlsx", ".xls"):
        df = pd.read_excel(path)
        return _table_to_text(df)
    if extension in (".doc", ".ppt"):
        raise PermanentError(f"Legacy format not supported: {extension}")

    return Path(path).read_text(encoding="utf-8", errors="ignore")


def _table_to_text(df: pd.DataFrame) -> str:
    rows: list[str] = []
    columns = list(df.columns)
    for _, row in df.iterrows():
        row_str = ", ".join(f"{col}: {row[col]}" for col in columns)
        rows.append(row_str)
    return "\n".join(rows)


def _chunk_text(text: str, target_tokens: int, overlap_tokens: int) -> list[Chunk]:
    tokenizer = _load_tokenizer()
    encoded = tokenizer(
        text,
        return_offsets_mapping=True,
        add_special_tokens=False,
    )
    input_ids = encoded.get("input_ids", [])
    offsets = encoded.get("offset_mapping", [])
    if not input_ids or not offsets:
        raise PermanentError("No tokens produced")

    step = max(1, target_tokens - overlap_tokens)
    chunks: list[Chunk] = []
    start = 0
    index = 0

    while start < len(input_ids):
        end = min(start + target_tokens, len(input_ids))
        if end <= start:
            break
        char_start = offsets[start][0]
        char_end = offsets[end - 1][1]
        if char_end <= char_start:
            start += step
            continue
        chunk_text = text[char_start:char_end]
        chunks.append(
            Chunk(
                index=index,
                text=chunk_text,
                char_start=char_start,
                char_end=char_end,
                token_start=start,
                token_end=end,
                token_count=end - start,
            )
        )
        index += 1
        start += step

    return chunks


def _embed_chunks(
    chunks: list[Chunk],
    tracker: CallTracker | None = None,
) -> list[list[float]]:
    if not chunks:
        return []
    embedder = get_text_embedder(768)
    batch_size = text_embed_batch_size()
    if tracker is not None:
        tracker.set_context(
            "text_embed",
            {
                "batch_size": batch_size,
                "backend": get_embedding_backend("text"),
            },
        )
    embeddings: list[list[float]] = []
    for start in range(0, len(chunks), batch_size):
        batch = [chunk.text for chunk in chunks[start : start + batch_size]]
        if tracker is None:
            batch_vectors = run_inference(
                "text",
                lambda batch=batch: embedder.encode(batch),
            )
        else:
            batch_vectors = timed_call(
                tracker,
                "text_embed",
                lambda batch=batch: run_inference(
                    "text",
                    lambda batch=batch: embedder.encode(batch),
                ),
            )
        if not batch_vectors:
            raise PermanentError("No embeddings produced")
        embeddings.extend(batch_vectors)
    if len(embeddings) != len(chunks):
        raise PermanentError("Embedding count mismatch")
    return embeddings


def ingest_document(
    *,
    source: IngestSource,
    config: Config,
    output_uri: str | None,
    pipeline_version: str,
    schema_version: str,
) -> PipelineResult:
    started_at = datetime.now(timezone.utc)
    extension = source.extension
    timer = StageTimer()
    calls = CallTracker()
    with timer.track("extract_text"):
        text = _extract_text(source.local_path, extension)
    if not text.strip() and config.enable_ocr and extension == ".pdf":
        with timer.track("ocr"):
            text = ocr_text_from_pdf(
                source.local_path,
                config.ocr_max_pages,
                base_uri=config.graph_root_uri(),
            )
    if not text.strip():
        raise PermanentError("No extractable text")

    with timer.track("chunk"):
        chunks = _chunk_text(
            text,
            config.chunk_target_tokens,
            config.chunk_overlap_tokens,
        )
    if not chunks:
        raise PermanentError("No chunks produced")

    with timer.track("embed"):
        embeddings = _embed_chunks(chunks, calls)

    output_root = output_uri or config.graph_root_uri()
    media_asset_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc)

    media_row = {
        "id": media_asset_id,
        "uri": source.uri,
        "media_type": "document",
        "content_type": source.content_type or "application/octet-stream",
        "size_bytes": source.size_bytes or 0,
        "source_bucket": source.bucket,
        "source_object": source.name,
        "source_generation": source.generation,
        "checksum": source.md5_hash or source.crc32c,
        "duration_ms": None,
        "width_px": None,
        "height_px": None,
        "frame_count": None,
        "sample_rate_hz": None,
        "channels": None,
        **tenancy_fields(
            org_id=source.org_id,
            site_id=source.site_id,
            stream_id=source.stream_id,
        ),
        "created_at": now,
        "pipeline_version": pipeline_version,
        "schema_version": schema_version,
    }

    chunk_core_rows = []
    chunk_text_rows = []
    chunk_vector_rows = []
    edge_rows = []

    for chunk, vector in zip(chunks, embeddings, strict=False):
        chunk_id = str(uuid.uuid4())
        chunk_core_rows.append(
            {
                "id": chunk_id,
                "media_asset_id": media_asset_id,
                "chunk_index": chunk.index,
                "char_start": chunk.char_start,
                "char_end": chunk.char_end,
                "token_start": chunk.token_start,
                "token_end": chunk.token_end,
                "token_count": chunk.token_count,
                "embedding_model": _pipeline_model(),
                **tenancy_fields(
                    org_id=source.org_id,
                    site_id=source.site_id,
                    stream_id=source.stream_id,
                ),
                "pipeline_version": pipeline_version,
                "schema_version": schema_version,
            }
        )
        chunk_text_rows.append({"content": chunk.text})
        chunk_vector_rows.append({"text_vector": vector})
        edge_rows.append(
            {
                "src_id": chunk_id,
                "dst_id": media_asset_id,
                "schema_version": schema_version,
            }
        )

    files: list[WriteResult] = []
    with timer.track("write_parquet"):
        compression = _doc_parquet_compression()
        row_group_size = _doc_parquet_row_group_size()
        jobs = [
            (
                [media_row],
                schema_for("MediaAsset", "core"),
                vertex_part_uri(
                    output_root, "MediaAsset", "core", str(uuid.uuid4())
                ),
            ),
            (
                chunk_core_rows,
                schema_for("DocChunk", "core"),
                vertex_part_uri(output_root, "DocChunk", "core", str(uuid.uuid4())),
            ),
            (
                chunk_text_rows,
                schema_for("DocChunk", "text"),
                vertex_part_uri(output_root, "DocChunk", "text", str(uuid.uuid4())),
            ),
            (
                chunk_vector_rows,
                schema_for("DocChunk", "vector"),
                vertex_part_uri(
                    output_root, "DocChunk", "vector", str(uuid.uuid4())
                ),
            ),
            (
                edge_rows,
                schema_for("DerivedFrom", "adj_list"),
                edge_part_uri(output_root, "DerivedFrom", str(uuid.uuid4())),
            ),
        ]
        files.extend(
            _write_parquet_parallel(
                jobs,
                compression=compression,
                row_group_size=row_group_size,
            )
        )

    parquet_bytes = sum(item.bytes_written for item in files)
    bytes_raw = source.size_bytes or 0
    token_total = sum(chunk.token_count for chunk in chunks)
    word_count = len(text.split())
    hashes: dict[str, str] = {}
    if source.content_hash_sha256:
        hashes["content_sha256"] = source.content_hash_sha256
    raw_timings_preview = timer.summary()
    stage_timings_preview = build_stage_timings(
        raw_timings_preview,
        {
            "extract_text": "decode_ms",
            "ocr": "decode_ms",
            "chunk": "finalize_ms",
            "embed": "embed_text_ms",
            "write_parquet": "write_parquet_ms",
            "write_manifest": "write_manifest_ms",
        },
    )
    manifest_metrics = manifest_metrics_subset(
        {
            "io": {
                "bytes_raw": bytes_raw,
                "bytes_parquet": parquet_bytes,
                "bytes_derived": parquet_bytes,
            },
            "quality": {
                "word_count": word_count,
                "token_count": token_total,
                "chunk_count": len(chunks),
            },
            "hashes": hashes,
            "embeddings": {
                "text": {
                    "count": len(chunks),
                    "dims": 768,
                }
            },
            "evidence": {
                "frames": 0,
                "snippets": len(chunks),
                "segments": 0,
            },
            "stage_timings_ms": stage_timings_preview,
        }
    )
    completed_at = datetime.now(timezone.utc)
    run_id = str(uuid.uuid4())
    with timer.track("write_manifest"):
        manifest = build_manifest(
            pipeline_version=pipeline_version,
            schema_version=schema_version,
            counts={
                "MediaAsset": 1,
                "DocChunk": len(chunk_core_rows),
                "DerivedFrom": len(edge_rows),
            },
            files=files,
            started_at=started_at,
            completed_at=completed_at,
            metrics=manifest_metrics,
        )
        manifest_path = manifest_uri(output_root, run_id)
        write_manifest(manifest, manifest_path, compact=True)
    raw_timings = timer.summary()
    stage_timings_ms = build_stage_timings(
        raw_timings,
        {
            "extract_text": "decode_ms",
            "ocr": "decode_ms",
            "chunk": "finalize_ms",
            "embed": "embed_text_ms",
            "write_parquet": "write_parquet_ms",
            "write_manifest": "write_manifest_ms",
        },
    )
    metrics = {
        "timings_ms": raw_timings,
        "stage_timings_ms": stage_timings_ms,
        "pipe_ms": round(sum(stage_timings_ms.values()), 2),
        "model_calls": calls.summary(),
        "io": {
            "bytes_raw": bytes_raw,
            "bytes_parquet": parquet_bytes,
            "bytes_derived": parquet_bytes,
        },
        "quality": {
            "word_count": word_count,
            "token_count": token_total,
            "chunk_count": len(chunks),
        },
        "hashes": hashes,
        "embeddings": {
            "text": {
                "count": len(chunks),
                "dims": 768,
            }
        },
        "evidence": {
            "frames": 0,
            "snippets": len(chunks),
            "segments": 0,
        },
    }

    return PipelineResult(
        counts={
            "MediaAsset": 1,
            "DocChunk": len(chunk_core_rows),
            "DerivedFrom": len(edge_rows),
        },
        manifest_uri=manifest_path,
        media_asset_id=media_asset_id,
        metrics=metrics,
    )
